import re
import os
import json
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

# ファイルパス設定
TEMPLATE_PATH = 'template.pptx'
OUTPUT_PATH = 'generated_context_map.pptx'
JSON_DATA_PATH = 'data.json'

print(f"現在の作業ディレクトリ: {os.getcwd()}")
print(f"テンプレートファイルが存在するか確認: {os.path.exists(TEMPLATE_PATH)}")

# 企業名を入力
company_name = input("企業名を入力してください: ")

# JSONデータを読み込む
try:
    with open(JSON_DATA_PATH, 'r', encoding='utf-8') as f:
        json_data = json.load(f)
    print(f"JSONデータを読み込みました: {len(json_data)}セクション")
except Exception as e:
    print(f"JSONデータの読み込みエラー: {e}")
    json_data = []

# テンプレートの文字マッピング
template_mapping = {
    "企業名": company_name,
    "企業名コンテクストマップ": f"{company_name}コンテクストマップ",  # 左上のタイトル用
    
    # 周囲の5つの大きなノード（section_name）
    "企業": "",
    "プロ": "",  
    "TOP": "",
    "業界": "",
    "グロ": "",
    
    # 各セクションの項目（四角形ノード）
    "あ": "", "い": "", "う": "", "え": "", "お": "",
    "か": "", "き": "", "く": "", "け": "", "こ": "",
    "さ": "", "し": "", "す": "", "せ": "", "そ": "",
    "た": "", "ち": "", "つ": "", "て": "", "と": "",
    "な": "", "に": "", "ぬ": "", "ね": "", "の": "",
    
    # ソーシャルトレンド（円形ノード）
    "一": "", "二": "", "三": "",
    "十": "", "十一": "", "十二": "",
    "四": "", "五": "", "六": "",
    "十三": "", "十四": "", "十五": "",
    "七": "", "八": "", "九": "",
}

# JSONデータからマッピングを更新
for idx, section in enumerate(json_data):
    code = section.get("code", "")
    section_name = section.get("section_name", "")
    items = section.get("items", [])
    social_trends = section.get("social_trend", [])
    
    # セクション名の配置
    if code == "A":
        template_mapping["企業"] = section_name
        item_keys = ["あ", "い", "う", "え", "お"]
        trend_keys = ["一", "二", "三"]
    elif code == "B":
        template_mapping["プロ"] = section_name
        item_keys = ["か", "き", "く", "け", "こ"]
        trend_keys = ["十", "十一", "十二"]
    elif code == "C":
        template_mapping["TOP"] = section_name
        item_keys = ["さ", "し", "す", "せ", "そ"]
        trend_keys = ["四", "五", "六"]
    elif code == "D":
        template_mapping["業界"] = section_name
        item_keys = ["た", "ち", "つ", "て", "と"]
        trend_keys = ["十三", "十四", "十五"]
    elif code == "E":
        template_mapping["グロ"] = section_name
        item_keys = ["な", "に", "ぬ", "ね", "の"]
        trend_keys = ["七", "八", "九"]
    else:
        continue
    
    # アイテムのタイトルを配置
    for i, item in enumerate(items[:5]):
        if i < len(item_keys):
            template_mapping[item_keys[i]] = item.get("title", "")
    
    # ソーシャルトレンドを配置
    for i, trend in enumerate(social_trends[:3]):
        if i < len(trend_keys):
            template_mapping[trend_keys[i]] = trend

# テキストの長さに応じて適切に改行を入れる関数
def format_text_for_shape(text, max_chars_per_line=7, is_circular=False, is_title=False):
    """テキストを適切な長さで改行する"""
    if not text:
        return text
    
    # タイトルの場合は改行しない
    if is_title:
        return text
    
    # シェイプのタイプに応じて改行基準を設定
    if is_circular:
        max_chars_per_line = 5
    else:
        max_chars_per_line = 7
    
    # 括弧がある場合は括弧で区切る
    if "（" in text and "）" in text:
        parts = text.split("（", 1)
        if len(parts) == 2:
            main_text = parts[0]
            bracket_text = "（" + parts[1]
            # メインテキストと括弧テキストをそれぞれ改行
            formatted_main = format_text_for_shape(main_text, max_chars_per_line, is_circular)
            return formatted_main + "\n" + bracket_text
    
    # 「・」で区切る
    if "・" in text and len(text) > 10:
        parts = text.split("・", 1)
        if len(parts) == 2:
            return parts[0] + "・\n" + parts[1]
    
    # 特定のキーワードでの改行
    keywords = ["による", "としての", "への", "での", "からの", "における"]
    for keyword in keywords:
        if keyword in text and len(text) > 12:
            parts = text.split(keyword, 1)
            if len(parts) == 2:
                return parts[0] + keyword + "\n" + parts[1]
    
    # 通常の改行処理
    if len(text) <= max_chars_per_line:
        return text
    
    result = []
    current_line = ""
    
    for i, char in enumerate(text):
        current_line += char
        # 次の文字が句読点でない場合に改行
        if len(current_line) >= max_chars_per_line and (i + 1 >= len(text) or text[i + 1] not in "、。"):
            result.append(current_line)
            current_line = ""
    
    if current_line:
        result.append(current_line)
    
    return "\n".join(result)

# フォントサイズを調整する関数
def adjust_font_size(shape, text, original_text, is_title=False):
    """テキストの長さに応じてフォントサイズを調整"""
    if not hasattr(shape, 'text_frame') or not shape.text_frame.paragraphs:
        return
    
    paragraph = shape.text_frame.paragraphs[0]
    
    # 元のフォントサイズを取得
    if paragraph.runs and paragraph.runs[0].font.size:
        original_size = paragraph.runs[0].font.size
    else:
        original_size = Pt(10)
    
    # タイトルの場合は大きめのフォントサイズを維持
    if is_title:
        new_size = max(Pt(12), original_size * 0.8)
    else:
        # テキストの行数とシェイプサイズを考慮して調整
        lines = text.count('\n') + 1
        text_length = len(text)
        
        if lines >= 3 or text_length > 25:
            new_size = max(Pt(7), original_size * 0.6)
        elif lines >= 2 or text_length > 15:
            new_size = max(Pt(8), original_size * 0.7)
        elif text_length > 10:
            new_size = max(Pt(9), original_size * 0.8)
        else:
            new_size = original_size * 0.9
    
    # フォントサイズを設定
    for run in paragraph.runs:
        run.font.size = new_size

try:
    # PowerPointテンプレートを開く
    prs = Presentation(TEMPLATE_PATH)
    print(f"\nテンプレートを読み込みました: スライド数 = {len(prs.slides)}")
    
    # スライドごとの処理
    total_replacements = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\nスライド {slide_idx+1} を処理中...")
        slide_replacements = 0
        
        # すべてのシェイプを処理
        for shape_idx, shape in enumerate(slide.shapes):
            # テキストフレームがあるか確認
            if not hasattr(shape, 'text_frame'):
                continue
            
            shape_text = shape.text_frame.text.strip()
            
            if not shape_text:
                continue
            
            # 左上のタイトルかどうかを判定
            is_title = (shape_idx == 0 and shape_text == "企業名コンテクストマップ")
            
            # マッピングに基づいて置換
            if shape_text in template_mapping and template_mapping[shape_text]:
                new_text = template_mapping[shape_text]
                
                # シェイプの形状を判定（円形かどうか）
                is_circular = False
                if hasattr(shape, 'width') and hasattr(shape, 'height'):
                    aspect_ratio = shape.width / shape.height if shape.height > 0 else 1
                    is_circular = 0.9 < aspect_ratio < 1.1
                
                # テキストのフォーマット
                formatted_text = format_text_for_shape(new_text, is_circular=is_circular, is_title=is_title)
                
                # テキストフレームの自動調整を設定
                text_frame = shape.text_frame
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                text_frame.word_wrap = True
                
                # 内部マージンを調整
                text_frame.margin_left = Pt(1)
                text_frame.margin_right = Pt(1)
                text_frame.margin_top = Pt(1)
                text_frame.margin_bottom = Pt(1)
                
                # 段落を設定
                if text_frame.paragraphs:
                    paragraph = text_frame.paragraphs[0]
                    paragraph.text = formatted_text
                    paragraph.alignment = PP_ALIGN.CENTER
                    
                    # フォントサイズを調整
                    adjust_font_size(shape, formatted_text, shape_text, is_title=is_title)
                else:
                    text_frame.text = formatted_text
                
                print(f"  シェイプ {shape_idx+1}: '{shape_text}' → '{new_text}'")
                slide_replacements += 1
        
        print(f"  スライド {slide_idx+1} での置換数: {slide_replacements}")
        total_replacements += slide_replacements
    
    print(f"\n総置換数: {total_replacements}")
    
    # 新しいファイルとして保存
    prs.save(OUTPUT_PATH)
    print(f"コンテクストマップを {OUTPUT_PATH} に保存しました。")
    
except Exception as e:
    print(f"エラーが発生しました: {e}")
    import traceback
    traceback.print_exc()